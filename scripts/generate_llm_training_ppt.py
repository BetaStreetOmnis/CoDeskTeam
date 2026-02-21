from __future__ import annotations

import argparse
from datetime import date
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ACCENT_BLUE = RGBColor(0x00, 0x50, 0xCA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xD9, 0xE0, 0xF3)
MUTED_TEXT = RGBColor(0xA9, 0xB4, 0xD6)
RED = RGBColor(0xFF, 0x4D, 0x4F)
GREEN = RGBColor(0x52, 0xC4, 0x1A)
CODE_BG = RGBColor(0x05, 0x12, 0x2B)


def _remove_slide(prs: Presentation, slide) -> None:
    slide_id = slide.slide_id
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld_id in list(sld_id_lst):
        if int(sld_id.get("id")) == slide_id:
            r_id = sld_id.get(qn("r:id"))
            if r_id:
                prs.part.drop_rel(r_id)
            sld_id_lst.remove(sld_id)
            return


def _clear_all_but_first_slide(prs: Presentation) -> None:
    for slide in list(prs.slides)[1:]:
        _remove_slide(prs, slide)


def _set_run_style(run, *, size: int, bold: bool = False, color: RGBColor = WHITE) -> None:
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color


def _add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.80), Inches(0.40), Inches(12.0), Inches(0.60))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _set_run_style(r, size=28, bold=True, color=WHITE)


def _add_subtitle(slide, subtitle: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.80), Inches(1.08), Inches(12.0), Inches(0.40))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    _set_run_style(r, size=14, bold=False, color=MUTED_TEXT)


def _add_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    heading: str,
    bullets: Iterable[str],
    heading_color: RGBColor = ACCENT_BLUE,
    bullet_color: RGBColor = LIGHT_TEXT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = heading
    _set_run_style(r0, size=16, bold=True, color=heading_color)

    for item in bullets:
        p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"• {item}"
        _set_run_style(r, size=12, bold=False, color=bullet_color)


def _add_two_col_list(
    slide,
    *,
    left_title: str,
    left_items: Iterable[str],
    right_title: str,
    right_items: Iterable[str],
) -> None:
    _add_card(
        slide,
        left=0.80,
        top=1.70,
        width=5.95,
        height=4.60,
        heading=left_title,
        bullets=left_items,
        heading_color=GREEN,
        bullet_color=LIGHT_TEXT,
    )
    _add_card(
        slide,
        left=6.95,
        top=1.70,
        width=5.95,
        height=4.60,
        heading=right_title,
        bullets=right_items,
        heading_color=RED,
        bullet_color=LIGHT_TEXT,
    )


def _add_simple_flow(slide) -> None:
    # Four-step flow: circles + arrows
    steps = [
        ("1", "选场景", "从高频、低风险的\n办公流程开始"),
        ("2", "给材料", "提供样例/数据/背景\n减少来回猜"),
        ("3", "迭代验证", "对照标准自检\n再交付"),
        ("4", "沉淀复用", "提示词库/知识库\n团队共享"),
    ]

    start_x = 1.05
    y = 2.10
    circle_w = 1.05
    gap = 1.80
    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (circle_w + gap)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(circle_w), Inches(circle_w)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT_BLUE
        circ.line.color.rgb = ACCENT_BLUE

        tf = circ.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        _set_run_style(r, size=22, bold=True, color=WHITE)

        label = slide.shapes.add_textbox(Inches(x - 0.25), Inches(y + 1.15), Inches(circle_w + 0.50), Inches(0.45))
        tf2 = label.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = title
        _set_run_style(r2, size=16, bold=True, color=WHITE)

        desc_box = slide.shapes.add_textbox(
            Inches(x - 0.45), Inches(y + 1.55), Inches(circle_w + 0.90), Inches(0.85)
        )
        tf3 = desc_box.text_frame
        tf3.clear()
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = desc
        _set_run_style(r3, size=12, bold=False, color=MUTED_TEXT)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + circle_w + 0.25), Inches(y + 0.15), Inches(1.10), Inches(0.75)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.color.rgb = ACCENT_BLUE


def _edit_cover_slide(cover_slide, *, company_name: str, training_date: str) -> None:
    # Main title box keeps typography: first line (blue), second line (white, large)
    for shape in cover_slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.name == "文本框 1":
            tf = shape.text_frame
            if len(tf.paragraphs) >= 2:
                p0 = tf.paragraphs[0]
                if p0.runs:
                    p0.runs[0].text = "生成式AI"
                p1 = tf.paragraphs[1]
                runs = p1.runs
                if len(runs) >= 2:
                    runs[0].text = "大模型"
                    runs[1].text = "培训"
                elif len(runs) == 1:
                    runs[0].text = "大模型培训"
            continue

        if "Copyright" in shape.text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = f"© {date.today().year} {company_name} · Internal Training Materials"
            _set_run_style(r, size=9, bold=False, color=MUTED_TEXT)
            continue

        if shape.text.strip() == "重庆旱獭信息技术有限公司":
            shape.text = company_name
            continue

        if "从“看清”到“读懂”" in shape.text:
            shape.text = "从“会问”到“会用”——让大模型成为你的工作助手"
            continue

    # Add date (small) on cover bottom-left if not present
    date_box = cover_slide.shapes.add_textbox(Inches(0.80), Inches(6.90), Inches(5.5), Inches(0.35))
    tf = date_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = training_date
    _set_run_style(r, size=12, bold=False, color=MUTED_TEXT)


def _add_llm_definition_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "大模型是什么（LLM）")
    _add_subtitle(slide, "一句话：基于海量数据训练的“通用语言能力引擎”，通过上下文生成/理解/推理文本。")

    left = slide.shapes.add_textbox(Inches(0.85), Inches(1.75), Inches(6.0), Inches(4.6))
    tf = left.text_frame
    tf.clear()
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "你需要记住的 3 个关键词"
    _set_run_style(r0, size=16, bold=True, color=ACCENT_BLUE)

    items = [
        "Token：模型处理文本的最小单位（字/词/片段）。",
        "上下文：你给的资料 + 对话历史，决定输出质量上限。",
        "幻觉：在信息不足时“看似合理”的编造，需要核对。",
    ]
    for it in items:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"• {it}"
        _set_run_style(r, size=12, bold=False, color=LIGHT_TEXT)

    # Simple IO diagram on the right
    box_w, box_h = Inches(1.8), Inches(0.95)
    arrow_w, arrow_h = Inches(0.40), Inches(0.55)
    x0, y0 = Inches(6.95), Inches(2.25)
    labels = ["输入\n(需求/材料)", "大模型\n(理解/生成)", "输出\n(草稿/方案)"]
    for i, label in enumerate(labels):
        x = x0 + i * (box_w + arrow_w)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, box_w, box_h)
        rect.fill.background()
        rect.line.color.rgb = ACCENT_BLUE
        rect.line.width = Pt(2)
        tfb = rect.text_frame
        tfb.clear()
        p = tfb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        _set_run_style(r, size=12, bold=True, color=WHITE)

        if i < len(labels) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w,
                y0 + Inches(0.20),
                arrow_w,
                arrow_h,
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.color.rgb = ACCENT_BLUE


def _add_prompt_template_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "提示词四件套：把AI当“同事”来写需求")
    _add_subtitle(slide, "写得越像需求文档，结果越稳定。推荐：角色｜目标｜材料｜输出格式｜约束。")

    _add_card(
        slide,
        left=0.85,
        top=1.75,
        width=5.9,
        height=4.7,
        heading="四件套（可记作 R-G-M-O）",
        bullets=[
            "Role：你希望它扮演谁（岗位/专家）",
            "Goal：要解决什么问题（可量化）",
            "Material：给足材料（原文/数据/样例）",
            "Output：指定格式（表格/要点/邮件/PPT大纲）",
        ],
    )

    code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.75), Inches(5.9), Inches(4.7))
    code.fill.solid()
    code.fill.fore_color.rgb = CODE_BG
    code.line.color.rgb = ACCENT_BLUE
    code.line.width = Pt(2)

    tf = code.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (
        "【角色】你是……\\n"
        "【目标】请在……内完成……\\n"
        "【材料】我提供：……（粘贴原文/数据）\\n"
        "【约束】不要编造；不确定就提问；标注引用。\\n"
        "【输出】请按：1) 结论 2) 依据 3) 下一步（表格）。"
    )
    f = r.font
    f.name = "Courier New"
    f.size = Pt(12)
    f.color.rgb = LIGHT_TEXT


def _add_use_cases_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "办公高频场景（可直接复用）")
    _add_subtitle(slide, "把“任务 + 材料 + 输出格式”写清楚，效率提升最明显。")

    cards = [
        ("写作润色", ["把这段话改成更专业/更口语", "限制：不改变事实、不夸大"]),
        ("会议纪要", ["把以下录音转写/要点整理", "输出：行动项/负责人/截止时间"]),
        ("资料总结", ["总结这份文档：结论+要点+风险", "输出：一页纸（500字内）"]),
        ("邮件/报告", ["根据要点写一封邮件", "语气：礼貌、简洁、可执行"]),
        ("表格分析", ["基于这份数据给洞察", "输出：3个发现+建议（表格）"]),
        ("PPT大纲", ["把主题拆成 8 页PPT 结构", "每页：标题+3要点+图示建议"]),
    ]

    left0, top0 = 0.85, 1.75
    col_w, row_h = 3.90, 2.15
    col_gap, row_gap = 0.33, 0.35

    for idx, (heading, lines) in enumerate(cards):
        r = idx // 3
        c = idx % 3
        x = left0 + c * (col_w + col_gap)
        y = top0 + r * (row_h + row_gap)
        _add_card(
            slide,
            left=x,
            top=y,
            width=col_w,
            height=row_h,
            heading=heading,
            bullets=lines,
            heading_color=ACCENT_BLUE,
            bullet_color=LIGHT_TEXT,
        )


def _add_department_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "部门场景案例（示例）")
    _add_subtitle(slide, "把“输入材料”准备好，比反复改提示词更重要。")

    cases = [
        ("销售", ["材料：客户背景/需求/竞品信息", "输出：电话话术+异议处理+跟进邮件"]),
        ("客服", ["材料：历史工单/FAQ/产品手册", "输出：标准回复+排障步骤+升级建议"]),
        ("人事/行政", ["材料：制度/流程/模板", "输出：公告文案+问答清单+培训题库"]),
        ("研发", ["材料：需求描述/日志/代码片段", "输出：排查思路+修复建议+测试用例"]),
    ]

    left0, top0 = 0.85, 1.75
    col_w, row_h = 6.0, 2.25
    col_gap, row_gap = 0.40, 0.35

    for idx, (heading, lines) in enumerate(cases):
        r = idx // 2
        c = idx % 2
        x = left0 + c * (col_w + col_gap)
        y = top0 + r * (row_h + row_gap)
        _add_card(
            slide,
            left=x,
            top=y,
            width=col_w,
            height=row_h,
            heading=heading,
            bullets=lines,
        )


def _add_safety_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "数据安全与合规：红线清单")
    _add_subtitle(slide, "原则：能不上传就不上传；必须上传先脱敏；输出用于对外前必须人工复核。")

    _add_two_col_list(
        slide,
        left_title="可以（建议）",
        left_items=[
            "公开信息、通用知识、已公开的政策与标准",
            "已脱敏/匿名化数据（去除姓名、电话、地址、账号等）",
            "公司模板/规范类文本（无敏感字段）",
            "对外材料的“润色/排版/摘要”（不改变事实）",
        ],
        right_title="不可以（禁止）",
        right_items=[
            "客户隐私、个人信息、未授权的身份证明材料",
            "商业机密：价格、合同、投标、财务、未发布路线图",
            "未公开源代码、密钥、账号口令、内部漏洞细节",
            "任何“让模型替你做最终决策/承诺”的内容",
        ],
    )


def _add_boundaries_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "能力与边界：别把AI当“事实源”")
    _add_subtitle(slide, "把它当“高效初稿生成器 + 思路陪练”，关键结论要用证据闭环。")

    cols = [
        ("擅长", ["起草与改写", "归纳总结", "生成结构化输出（表格/清单）", "提供多方案与权衡建议"]),
        ("风险点", ["信息不足时会编造（幻觉）", "可能遗漏细节/误读上下文", "对公司内部知识不了解（除非提供材料）", "专业/法律结论不具备责任主体"]),
        ("你要做", ["给材料：原文/数据/样例", "定标准：输出格式+验收口径", "做复核：关键事实/数字/引用", "做留痕：保存输入与版本（便于追溯）"]),
    ]

    left0, top0 = 0.85, 1.75
    col_w, col_gap = 3.90, 0.30
    for i, (h, b) in enumerate(cols):
        _add_card(
            slide,
            left=left0 + i * (col_w + col_gap),
            top=top0,
            width=col_w,
            height=4.8,
            heading=h,
            bullets=b,
            heading_color=ACCENT_BLUE if h != "风险点" else RED,
        )


def _add_goals_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "本次培训：目标与结构")
    _add_subtitle(slide, "适用对象：全员｜建议时长：60 分钟｜目标：安全、稳定、可复用地用AI提效。")

    cards = [
        ("理解", ["大模型是什么/能做什么", "为什么会“胡说”", "如何降低风险"]),
        ("会用", ["提示词模板与迭代方法", "办公高频场景示例", "让输出更稳的检查清单"]),
        ("能落地", ["部门场景选题方法", "提示词库/知识库沉淀", "安全合规与治理要点"]),
    ]
    left0, top0 = 0.85, 1.75
    col_w, col_gap = 3.90, 0.30
    for i, (h, b) in enumerate(cards):
        _add_card(
            slide,
            left=left0 + i * (col_w + col_gap),
            top=top0,
            width=col_w,
            height=4.8,
            heading=h,
            bullets=b,
        )


def _add_workflow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "落地方法：从个人提效到团队能力")
    _add_subtitle(slide, "推荐闭环：场景 → 材料 → 验证 → 沉淀（提示词库/知识库/流程）。")
    _add_simple_flow(slide)

    footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.25), Inches(12.0), Inches(0.55))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "建议产出：部门 Top10 提效场景清单 + 3 个可复用提示词模板 + 1 个共享知识库入口"
    _set_run_style(r, size=13, bold=False, color=LIGHT_TEXT)


def _add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.30), Inches(12.0), Inches(1.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Q&A"
    _set_run_style(r, size=54, bold=True, color=WHITE)

    box2 = slide.shapes.add_textbox(Inches(0.85), Inches(3.60), Inches(12.0), Inches(0.8))
    tf2 = box2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "谢谢｜欢迎提问与分享你的场景"
    _set_run_style(r2, size=20, bold=True, color=ACCENT_BLUE)


def build_ppt(*, template_path: Path, output_path: Path, company_name: str, training_date: str) -> None:
    prs = Presentation(str(template_path))
    if not prs.slides:
        raise RuntimeError("Template PPTX has no slides; cannot reuse its style.")

    cover = prs.slides[0]
    _edit_cover_slide(cover, company_name=company_name, training_date=training_date)
    _clear_all_but_first_slide(prs)

    _add_goals_slide(prs)
    _add_llm_definition_slide(prs)
    _add_boundaries_slide(prs)
    _add_prompt_template_slide(prs)
    _add_use_cases_slide(prs)
    _add_department_slide(prs)
    _add_safety_slide(prs)
    _add_workflow_slide(prs)
    _add_closing_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a company-wide LLM training PPTX with JetLinks-style template.")
    parser.add_argument(
        "--template",
        default="/Users/chenhao/Library/Containers/com.tencent.xinWeChat/Data/Documents/xwechat_files/wxid_or5tukb9c7mi21_2645/temp/drag/JetLinks AI视觉分析一体机2026版V1.2(1).pptx",
        help="Template PPTX path.",
    )
    parser.add_argument(
        "--output",
        default=str(Path("outputs") / f"大模型培训_全员_{date.today().isoformat()}.pptx"),
        help="Output PPTX path.",
    )
    parser.add_argument("--company", default="（公司名称）", help="Company name shown on cover.")
    parser.add_argument("--date", dest="training_date", default=date.today().isoformat(), help="Training date shown on cover.")

    args = parser.parse_args()
    build_ppt(
        template_path=Path(args.template),
        output_path=Path(args.output),
        company_name=args.company,
        training_date=args.training_date,
    )
    print(f"OK: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
