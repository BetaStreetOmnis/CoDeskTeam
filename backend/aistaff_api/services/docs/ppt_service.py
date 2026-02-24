from __future__ import annotations

from datetime import datetime
import os
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
from uuid import uuid4

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ...config import Settings
from ...output_cleanup import maybe_cleanup_outputs_dir
from ...url_utils import abs_url
from ..auth_service import create_download_token


class PptDocService:
    def __init__(self, settings: Settings) -> None:
        self._settings = settings
        self._settings.outputs_dir.mkdir(parents=True, exist_ok=True)
        maybe_cleanup_outputs_dir(
            self._settings.outputs_dir,
            ttl_seconds=max(0, int(self._settings.outputs_ttl_hours)) * 3600,
        )

    def _resolve_soffice_bin(self) -> str | None:
        configured = str(os.getenv("AISTAFF_SOFFICE_BIN") or "").strip()
        candidates = [
            configured,
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "soffice",
            "libreoffice",
        ]
        seen: set[str] = set()
        for raw in candidates:
            value = str(raw or "").strip()
            if not value or value in seen:
                continue
            seen.add(value)

            if os.path.isabs(value):
                path = Path(value).expanduser()
                if path.exists() and path.is_file():
                    return str(path)
                continue

            resolved = shutil.which(value)
            if resolved:
                return resolved
        return None

    def _build_soffice_env(self, tmp_dir: Path) -> tuple[dict, str]:
        profile_dir = (tmp_dir / "lo_profile").resolve()
        profile_dir.mkdir(parents=True, exist_ok=True)
        env = os.environ.copy()
        for key in ("HOME", "USERPROFILE", "TMPDIR", "TEMP", "TMP"):
            env[key] = str(tmp_dir)
        env.setdefault("LANG", "en_US.UTF-8")
        env.setdefault("LC_ALL", "en_US.UTF-8")
        return env, f"-env:UserInstallation={profile_dir.as_uri()}"

    def _soffice_convert(self, *, soffice: str, tmp_dir: Path, out_dir: Path, input_path: Path, fmt: str) -> None:
        env, user_install = self._build_soffice_env(tmp_dir)
        cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--nofirststartwizard",
            "--invisible",
            user_install,
            "--convert-to",
            fmt,
            "--outdir",
            str(out_dir),
            str(input_path),
        ]
        subprocess.run(cmd, capture_output=True, check=False, timeout=60, env=env)

    def _create_ppt_cover_preview(self, ppt_path: Path) -> dict | None:
        soffice = self._resolve_soffice_bin()
        if not soffice:
            return None

        tmp_dir = Path(tempfile.mkdtemp(prefix="aistaff-ppt-preview-"))
        try:
            self._soffice_convert(
                soffice=soffice,
                tmp_dir=tmp_dir,
                out_dir=tmp_dir,
                input_path=ppt_path,
                fmt="png:impress_png_Export",
            )

            png_path = tmp_dir / f"{ppt_path.stem}.png"
            if not png_path.exists():
                png_candidates = sorted(tmp_dir.glob("*.png"), key=lambda item: item.stat().st_mtime, reverse=True)
                if not png_candidates:
                    return None
                png_path = png_candidates[0]

            preview_file_id = f"{uuid4().hex}.png"
            preview_path = (self._settings.outputs_dir / preview_file_id).resolve()
            preview_path.write_bytes(png_path.read_bytes())

            token = create_download_token(settings=self._settings, file_id=preview_file_id)
            return {
                "preview_image_file_id": preview_file_id,
                "preview_image_url": abs_url(self._settings, f"/api/files/{preview_file_id}?token={token}"),
            }
        except Exception:
            return None
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    def _normalize_ppt_style(self, value: str | None) -> str:
        raw = str(value or "").strip().lower()
        if not raw:
            return "auto"
        alias = {
            "auto": "auto",
            "modern": "modern_blue",
            "modern_blue": "modern_blue",
            "blue": "modern_blue",
            "minimal": "minimal_gray",
            "minimal_gray": "minimal_gray",
            "gray": "minimal_gray",
            "dark": "dark_tech",
            "dark_tech": "dark_tech",
            "tech": "dark_tech",
            "warm": "warm_business",
            "warm_business": "warm_business",
            "business": "warm_business",
            "template_jetlinks": "template_jetlinks",
            "jetlinks": "template_jetlinks",
            "template_team": "template_team",
            "team": "template_team",
            "team_style": "template_team",
            "企业蓝": "template_team",
            "科技蓝": "template_jetlinks",
            "深色科技": "dark_tech",
            "暖色商务": "warm_business",
            "极简灰": "minimal_gray",
            "现代蓝": "modern_blue",
        }
        return alias.get(raw, "auto")

    def _normalize_ppt_layout_mode(self, value: str | None) -> str:
        raw = str(value or "").strip().lower()
        if not raw:
            return "auto"
        alias = {
            "auto": "auto",
            "focus": "focus",
            "single": "single_column",
            "single_column": "single_column",
            "two": "two_column",
            "two_column": "two_column",
            "double": "two_column",
            "cards": "cards",
            "card_grid": "cards",
            "聚焦": "focus",
            "单栏": "single_column",
            "双栏": "two_column",
            "卡片": "cards",
        }
        return alias.get(raw, "auto")

    _FILE_ID_RE = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,199}$")

    def _safe_template_from_outputs(self, file_id: str | None) -> Path | None:
        fid = str(file_id or "").strip()
        if not fid or ".." in fid or not self._FILE_ID_RE.match(fid):
            return None
        base = self._settings.outputs_dir.resolve()
        full = (base / fid).resolve()
        if full == base or not str(full).startswith(str(base) + os.sep):
            return None
        if not full.exists() or not full.is_file():
            return None
        if full.suffix.lower() != ".pptx":
            return None
        return full

    def _ppt_style_profiles(self) -> dict[str, dict[str, str]]:
        return {
            "modern_blue": {
                "cover_from": "0F172A",
                "cover_to": "1E3A8A",
                "cover_text": "FFFFFF",
                "primary": "0F172A",
                "accent": "4F46E5",
                "accent2": "0EA5E9",
                "text": "111827",
                "muted": "64748B",
                "border": "E2E8F0",
                "bg": "F8FAFC",
                "surface": "FFFFFF",
                "chip_bg": "E0E7FF",
                "success": "14B8A6",
            },
            "minimal_gray": {
                "cover_from": "1F2937",
                "cover_to": "374151",
                "cover_text": "FFFFFF",
                "primary": "111827",
                "accent": "475569",
                "accent2": "64748B",
                "text": "1F2937",
                "muted": "6B7280",
                "border": "D1D5DB",
                "bg": "F9FAFB",
                "surface": "FFFFFF",
                "chip_bg": "E5E7EB",
                "success": "0F766E",
            },
            "dark_tech": {
                "cover_from": "020617",
                "cover_to": "111827",
                "cover_text": "E2E8F0",
                "primary": "E5E7EB",
                "accent": "22D3EE",
                "accent2": "6366F1",
                "text": "E2E8F0",
                "muted": "94A3B8",
                "border": "334155",
                "bg": "0B1220",
                "surface": "111827",
                "chip_bg": "1E293B",
                "success": "34D399",
            },
            "warm_business": {
                "cover_from": "422006",
                "cover_to": "7C2D12",
                "cover_text": "FFFFFF",
                "primary": "3F2A1D",
                "accent": "C2410C",
                "accent2": "F59E0B",
                "text": "3F2A1D",
                "muted": "7C5E49",
                "border": "E7D8C8",
                "bg": "FFF7ED",
                "surface": "FFFFFF",
                "chip_bg": "FDE68A",
                "success": "059669",
            },
        }

    def _ppt_template_candidates(self, style: str) -> list[Path]:
        base = self._settings.app_root / "backend" / "aistaff_api" / "assets" / "ppt_templates"
        jetlinks = base / "jetlinks_ai_vision_template.pptx"
        team = base / "team_style_template.pptx"

        values: list[Path] = []
        env_path = str(os.getenv("AISTAFF_PPT_TEMPLATE") or "").strip()
        if env_path:
            values.append(Path(env_path).expanduser())

        if style == "template_jetlinks":
            values.append(jetlinks)
        elif style == "template_team":
            values.append(team)
        elif style == "auto":
            # Prefer the simpler team template by default. If it can't satisfy the requested slide count,
            # we fall back to the richer JetLinks template.
            values.extend([team, jetlinks])

        seen: set[str] = set()
        out: list[Path] = []
        for p in values:
            key = str(p)
            if key in seen:
                continue
            seen.add(key)
            out.append(p)
        return out

    async def create_pptx(
        self,
        *,
        title: str,
        slides: list[dict],
        style: str | None = None,
        layout_mode: str | None = None,
        template_file_id: str | None = None,
        template_mode: str | None = None,
        template_keep_images: bool | None = None,
        template_content_indices: list[int] | None = None,
    ) -> dict:
        normalized_style = self._normalize_ppt_style(style)
        normalized_layout_mode = self._normalize_ppt_layout_mode(layout_mode)
        style_profiles = self._ppt_style_profiles()
        effective_style = normalized_style if normalized_style in style_profiles else "modern_blue"
        template_failures: list[dict[str, str]] = []
        strict_template = str(os.getenv("AISTAFF_PPT_TEMPLATE_STRICT") or "").strip().lower() in {
            "1",
            "true",
            "yes",
            "on",
        }
        template_override = self._safe_template_from_outputs(template_file_id)

        def _effective_keep_images() -> bool:
            if template_keep_images is not None:
                return bool(template_keep_images)
            raw = str(os.getenv("AISTAFF_PPT_KEEP_TEMPLATE_IMAGES") or "").strip().lower()
            if raw:
                return raw in {"1", "true", "yes", "on"}
            return True

        def _effective_template_mode() -> str:
            raw = str(template_mode or os.getenv("AISTAFF_PPT_TEMPLATE_MODE") or "").strip().lower()
            return raw or "reuse"

        def _build_result(path: Path) -> dict:
            file_id = path.name
            token = create_download_token(settings=self._settings, file_id=file_id)
            result = {
                "file_id": file_id,
                "filename": file_id,
                "download_url": abs_url(self._settings, f"/api/files/{file_id}?token={token}"),
                "style": normalized_style,
                "layout_mode": normalized_layout_mode,
            }
            preview_meta = self._create_ppt_cover_preview(path)
            if preview_meta:
                result.update(preview_meta)
            return result

        def _normalize_template_indices(indices: list[int] | None, total: int) -> list[int]:
            out: list[int] = []
            for raw in indices or []:
                try:
                    idx = int(raw)
                except Exception:
                    continue
                if idx == 1:
                    continue
                if idx < 1 or idx > total:
                    continue
                if idx in out:
                    continue
                out.append(idx)
            return out

        def _rgb(hex_value: str) -> RGBColor:
            value = (hex_value or "").strip().lstrip("#")
            if len(value) != 6:
                return RGBColor(0, 0, 0)
            return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

        def _normalize_text(value: object, *, max_chars: int, fallback: str = "") -> str:
            text_value = " ".join(str(value or "").split())
            if not text_value:
                return fallback
            if max_chars <= 1:
                return text_value[: max(0, max_chars)]
            if len(text_value) <= max_chars:
                return text_value
            return f"{text_value[: max_chars - 1]}…"

        bullet_prefix_re = re.compile(r"^\s*(?:[-*•·]+|\d+[.)、])\s*")

        def _normalize_bullet(value: object) -> str:
            text_value = _normalize_text(value, max_chars=140)
            if not text_value:
                return ""
            text_value = bullet_prefix_re.sub("", text_value).strip()
            # Keep bullets compact so the template/programmatic layouts don't look crowded.
            return _normalize_text(text_value, max_chars=56)

        def _cover_title_lines(value: str) -> list[str]:
            text_value = _normalize_text(value, max_chars=52, fallback="演示文稿")
            if len(text_value) <= 16:
                return [text_value]

            split_pos = -1
            mid = len(text_value) // 2
            for marker in ["：", ":", "—", "-", "·", " "]:
                idx = text_value.find(marker, max(6, mid - 8), min(len(text_value) - 6, mid + 8))
                if idx > 0:
                    split_pos = idx + (0 if marker == " " else 1)
                    break

            if split_pos < 0:
                split_pos = mid

            left = text_value[:split_pos].strip()
            right = text_value[split_pos:].strip()
            if not right:
                return [text_value]
            return [left, right]

        def _cover_title_size(text_len: int) -> Pt:
            if text_len <= 14:
                return Pt(52)
            if text_len <= 20:
                return Pt(44)
            if text_len <= 30:
                return Pt(38)
            if text_len <= 40:
                return Pt(32)
            return Pt(28)

        def _title_size(text_len: int) -> Pt:
            if text_len <= 14:
                return Pt(34)
            if text_len <= 22:
                return Pt(30)
            if text_len <= 30:
                return Pt(27)
            return Pt(24)

        def _font_size_for_bullets(items: list[str]) -> Pt:
            count = len(items)
            longest = max((len(item) for item in items), default=0)
            total = sum(len(item) for item in items)
            if count <= 4 and longest <= 24 and total <= 96:
                return Pt(24)
            if count <= 5 and longest <= 32 and total <= 145:
                return Pt(21)
            if count <= 6 and longest <= 42 and total <= 220:
                return Pt(19)
            return Pt(17)

        def _split_bullets(items: list[str]) -> tuple[list[str], list[str]]:
            left: list[str] = []
            right: list[str] = []
            left_weight = 0
            right_weight = 0
            for item in items:
                weight = max(1, len(item))
                if left_weight <= right_weight:
                    left.append(item)
                    left_weight += weight
                else:
                    right.append(item)
                    right_weight += weight
            return left, right

        normalized_title = _normalize_text(title, max_chars=56, fallback="演示文稿")

        normalized_slides: list[dict[str, object]] = []
        for idx, raw_slide in enumerate(slides, 1):
            source = raw_slide or {}
            slide_title = _normalize_text(source.get("title"), max_chars=34, fallback=f"第 {idx} 页")
            raw_bullets = source.get("bullets") or []
            normalized_bullets = [_normalize_bullet(item) for item in raw_bullets]
            cleaned_bullets = [item for item in normalized_bullets if item]
            if len(cleaned_bullets) > 8:
                cleaned_bullets = [*cleaned_bullets[:7], "其余细节可在讲解时展开说明。"]
            if not cleaned_bullets:
                cleaned_bullets = ["结合业务现状补充关键结论与下一步行动。"]
            normalized_slides.append({"title": slide_title, "bullets": cleaned_bullets})

        cover_topics = [str(item.get("title") or "").strip() for item in normalized_slides[:4]]
        cover_topics = [item for item in cover_topics if item]
        cover_subtitle = " · ".join(cover_topics) if cover_topics else "通识入门 · 业务落地 · 实践复盘"
        if len(normalized_slides) > 4 and cover_topics:
            cover_subtitle = f"{cover_subtitle} · …"
        cover_subtitle = _normalize_text(cover_subtitle, max_chars=62, fallback="通识入门 · 业务落地 · 实践复盘")

        def _shape_text(shape) -> str:  # noqa: ANN001
            if not getattr(shape, "has_text_frame", False):
                return ""
            return "\n".join((shape.text or "").splitlines()).strip()

        def _shape_font_score(shape) -> float:  # noqa: ANN001
            if not getattr(shape, "has_text_frame", False):
                return 0.0
            best = 0.0
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not str(run.text or "").strip():
                            continue
                        if run.font.size is not None:
                            best = max(best, float(run.font.size.pt))
            except Exception:
                return best
            return best

        def _set_shape_text(shape, value: str) -> None:  # noqa: ANN001
            if not getattr(shape, "has_text_frame", False):
                return
            try:
                shape.text = str(value or "")
            except Exception:
                return

        def _remove_shape(shape) -> None:  # noqa: ANN001
            try:
                element = shape._element
                element.getparent().remove(element)
            except Exception:
                return

        def _iter_shapes(items):  # noqa: ANN001
            for shape in items:
                yield shape
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for child in _iter_shapes(shape.shapes):
                            yield child
                    except Exception:
                        continue

        def _shape_area(shape) -> int:  # noqa: ANN001
            try:
                return int(shape.width) * int(shape.height)
            except Exception:
                return 0

        def _is_background_shape(shape, slide_width: int, slide_height: int) -> bool:  # noqa: ANN001
            try:
                width = int(shape.width)
                height = int(shape.height)
            except Exception:
                return False
            if width <= 0 or height <= 0:
                return False
            return width >= int(slide_width * 0.85) and height >= int(slide_height * 0.85)

        def _is_background_picture(shape, slide_width: int, slide_height: int) -> bool:  # noqa: ANN001
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                return False
            return _is_background_shape(shape, slide_width, slide_height)

        def _cleanup_slide(slide, slide_width: int, slide_height: int, *, remove_pictures: bool = True) -> None:  # noqa: ANN001
            for shape in list(slide.shapes):
                if _is_background_shape(shape, slide_width, slide_height):
                    continue
                shape_type = getattr(shape, "shape_type", None)
                if shape_type in {MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE}:
                    _remove_shape(shape)
                    continue
                if shape_type == MSO_SHAPE_TYPE.GROUP:
                    _remove_shape(shape)
                    continue
                if shape_type == MSO_SHAPE_TYPE.PICTURE and remove_pictures:
                    _remove_shape(shape)
                    continue
                if remove_pictures and not getattr(shape, "has_text_frame", False):
                    _remove_shape(shape)
                    continue
            for shape in _iter_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False):
                    _set_shape_text(shape, "")

        def _pick_title_shape(slide) -> object | None:  # noqa: ANN001
            text_shapes = [shape for shape in _iter_shapes(slide.shapes) if getattr(shape, "has_text_frame", False)]
            if not text_shapes:
                return None
            for shape in text_shapes:
                if getattr(shape, "is_placeholder", False):
                    try:
                        if int(shape.placeholder_format.type) in {int(PP_PLACEHOLDER.TITLE), int(PP_PLACEHOLDER.CENTER_TITLE)}:
                            return shape
                    except Exception:
                        continue

            def _is_body_placeholder(item) -> bool:  # noqa: ANN001
                if not getattr(item, "is_placeholder", False):
                    return False
                try:
                    return int(item.placeholder_format.type) in {int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)}
                except Exception:
                    return False

            top_band = [shape for shape in text_shapes if int(shape.top) <= int(Inches(1.8))]
            candidates = top_band or text_shapes
            non_body = [shape for shape in candidates if not _is_body_placeholder(shape)]
            if non_body:
                candidates = non_body

            def _score(item) -> tuple[float, int, int]:  # noqa: ANN001
                try:
                    top = int(item.top)
                except Exception:
                    top = 0
                return (_shape_font_score(item), _shape_area(item), -top)

            return max(candidates, key=_score)

        def _pick_body_shape(slide, title_shape) -> object | None:  # noqa: ANN001
            text_shapes = [shape for shape in _iter_shapes(slide.shapes) if getattr(shape, "has_text_frame", False)]
            if not text_shapes:
                return None
            candidates = [shape for shape in text_shapes if shape is not title_shape]
            if title_shape is not None:
                candidates = [
                    shape
                    for shape in candidates
                    if int(shape.top) >= int(title_shape.top) + int(Inches(0.4))
                ]
            if not candidates:
                return None
            return max(candidates, key=_shape_area)

        def _find_blank_layout(prs_obj: Presentation):  # noqa: ANN001
            layouts = list(prs_obj.slide_layouts)
            if not layouts:
                return None
            return min(layouts, key=lambda layout: len(layout.shapes))

        def _duplicate_slide(prs_obj: Presentation, source_slide, blank_layout) -> object:  # noqa: ANN001
            layout = blank_layout or source_slide.slide_layout
            new_slide = prs_obj.slides.add_slide(layout)
            # remove any auto-added shapes from layout
            for shape in list(new_slide.shapes):
                _remove_shape(shape)
            for shape in source_slide.shapes:
                try:
                    new_shape = shape._element.clone()
                except Exception:
                    from copy import deepcopy

                    new_shape = deepcopy(shape._element)
                new_slide.shapes._spTree.insert_element_before(new_shape, "p:extLst")
            return new_slide

        def _pick_base_content_slide_index(prs_obj: Presentation) -> int:
            slides = list(prs_obj.slides)
            if len(slides) <= 1:
                return 0
            raw_override = str(os.getenv("AISTAFF_PPT_TEMPLATE_CONTENT_INDEX") or "").strip()
            if raw_override.isdigit():
                override = int(raw_override)
                if 1 <= override <= len(slides):
                    return max(0, override - 1)
            candidates: list[tuple[tuple[int, int, int, int], int]] = []
            for idx, slide in enumerate(slides):
                if idx == 0:
                    continue
                shape_count = len(slide.shapes)
                text_count = sum(1 for s in slide.shapes if getattr(s, "has_text_frame", False))
                group_count = sum(1 for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.GROUP)
                pic_count = sum(1 for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE)
                if text_count == 0:
                    continue
                score = (group_count, pic_count, shape_count, -text_count)
                candidates.append((score, idx))
            if not candidates:
                return 1
            candidates.sort(key=lambda item: item[0])
            return candidates[0][1]

        def _parse_template_indices(raw: str, total: int) -> list[int]:
            tokens = [item.strip() for item in str(raw or "").split(",") if item.strip()]
            out: list[int] = []
            for token in tokens:
                if not token.isdigit():
                    continue
                idx = int(token)
                if 1 <= idx <= total and idx not in out:
                    out.append(idx)
            return out

        def _split_groups(items: list[str], groups: int) -> list[list[str]]:
            if groups <= 0:
                return []
            out = [[] for _ in range(groups)]
            weights = [0 for _ in range(groups)]
            for item in items:
                idx = min(range(groups), key=lambda i: weights[i])
                out[idx].append(item)
                weights[idx] += max(1, len(item))
            return out

        def _remove_slide(prs_obj: Presentation, index: int) -> None:
            slide_id_list = prs_obj.slides._sldIdLst
            slide_ids = list(slide_id_list)
            if index < 0 or index >= len(slide_ids):
                return
            slide_id = slide_ids[index]
            rel_id = slide_id.rId
            prs_obj.part.drop_rel(rel_id)
            slide_id_list.remove(slide_id)

        def _fill_cover_slide(slide) -> None:  # noqa: ANN001
            text_shapes = [shape for shape in _iter_shapes(slide.shapes) if getattr(shape, "has_text_frame", False)]
            if not text_shapes:
                return

            title_shape = max(text_shapes, key=_shape_font_score)
            title_lines = _cover_title_lines(normalized_title)

            date_shape = None
            subtitle_shape = None
            for shape in text_shapes:
                if shape is title_shape:
                    continue
                value = _shape_text(shape)
                if not value:
                    continue
                if re.search(r"\d{4}[-/.]\d{1,2}[-/.]\d{1,2}", value):
                    date_shape = shape
                    continue
                if "©" in value or "copyright" in value.lower():
                    continue
                if shape.top <= title_shape.top:
                    continue
                if subtitle_shape is None or int(shape.width) > int(subtitle_shape.width):
                    subtitle_shape = shape

            for shape in text_shapes:
                _set_shape_text(shape, "")

            _set_shape_text(title_shape, "\n".join(title_lines))
            if subtitle_shape is not None:
                _set_shape_text(subtitle_shape, cover_subtitle)
            if date_shape is not None:
                _set_shape_text(date_shape, datetime.now().strftime("%Y-%m-%d"))

        def _fill_clean_content_slide(slide, slide_title: str, bullets: list[str]) -> None:  # noqa: ANN001
            title_shape = _pick_title_shape(slide)
            body_shape = _pick_body_shape(slide, title_shape)

            if title_shape is None:
                title_shape = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(11.2), Inches(0.8))
            _set_shape_text(title_shape, _normalize_text(slide_title, max_chars=34, fallback="—"))

            if body_shape is None:
                body_shape = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.2), Inches(4.8))

            tf = body_shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP

            items = bullets or ["结合业务场景补充关键结论与下一步行动。"]
            font_size = _font_size_for_bullets(items)
            for idx, item in enumerate(items[:8]):
                paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                paragraph.text = f"• {item}"
                paragraph.font.name = "微软雅黑"
                paragraph.font.size = font_size

        def _pick_content_shapes(text_boxes: list, *, title_shape, subtitle_shape) -> list:  # noqa: ANN001
            min_width = int(Inches(2.4))
            min_height = int(Inches(0.45))
            candidates = []
            allowed_types = {MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER}

            for shape in text_boxes:
                if shape is title_shape or shape is subtitle_shape:
                    continue
                if getattr(shape, "shape_type", None) not in allowed_types:
                    continue

                value = _shape_text(shape)
                if not value:
                    continue

                normalized = " ".join(value.split())
                semantic = re.sub(r"[^0-9A-Za-z一-鿿]", "", normalized)
                if len(semantic) < 6:
                    continue
                if int(shape.width) < min_width or int(shape.height) < min_height:
                    continue

                score = (int(shape.width) * int(shape.height)) + (len(normalized) * 5000)
                candidates.append((score, shape))

            if not candidates:
                return []

            candidates.sort(key=lambda item: item[0], reverse=True)
            return [shape for _, shape in candidates[:2]]

        def _fill_content_slide(slide, slide_title: str, bullets: list[str]) -> None:  # noqa: ANN001
            text_boxes = [
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and _shape_text(shape)
            ]
            text_boxes.sort(key=lambda shape: (int(shape.top), int(shape.left)))
            if not text_boxes:
                return

            title_shape = text_boxes[0]
            _set_shape_text(title_shape, _normalize_text(slide_title, max_chars=34, fallback="—"))

            subtitle_shape = None
            for shape in text_boxes[1:]:
                if int(shape.top) - int(title_shape.top) > int(Inches(1.9)):
                    break
                if getattr(shape, "shape_type", None) not in {MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER}:
                    continue
                if int(shape.width) >= int(Inches(2.6)):
                    subtitle_shape = shape
                    break

            if subtitle_shape is not None:
                sub = " · ".join(bullets[:2]) if bullets else "结合业务场景补充关键结论"
                _set_shape_text(subtitle_shape, _normalize_text(sub, max_chars=78, fallback="结合业务场景补充关键结论"))

            content_shapes = _pick_content_shapes(text_boxes, title_shape=title_shape, subtitle_shape=subtitle_shape)
            if not content_shapes:
                inject = slide.shapes.add_textbox(Inches(0.9), Inches(5.78), Inches(11.5), Inches(1.25))
                inject_tf = inject.text_frame
                inject_tf.clear()
                inject_tf.word_wrap = True
                for idx, item in enumerate((bullets or ["结合业务场景补充要点"])[:3]):
                    paragraph = inject_tf.paragraphs[0] if idx == 0 else inject_tf.add_paragraph()
                    paragraph.text = f"• {item}"
                    paragraph.font.name = "微软雅黑"
                    paragraph.font.size = Pt(15)
                return

            groups = _split_groups(bullets, len(content_shapes))
            for idx, shape in enumerate(content_shapes):
                group = groups[idx] if idx < len(groups) else []
                if not group:
                    _set_shape_text(shape, "• 结合业务场景补充要点")
                    continue
                lines = [f"• {item}" for item in group[:6]]
                _set_shape_text(shape, "\n".join(lines))

        def _render_with_template(template_path: Path) -> dict:
            template = template_path.expanduser().resolve()
            prs_t = Presentation(str(template))
            if len(prs_t.slides) < 1:
                raise ValueError("template must include at least 1 slide")

            needed = len(normalized_slides)
            if needed <= 0:
                raise ValueError("no slides to render")

            slide_width = int(prs_t.slide_width)
            slide_height = int(prs_t.slide_height)
            blank_layout = _find_blank_layout(prs_t)

            def _shape_ph_type(shape) -> int | None:  # noqa: ANN001
                if not getattr(shape, "is_placeholder", False):
                    return None
                try:
                    return int(shape.placeholder_format.type)  # type: ignore[arg-type]
                except Exception:
                    return None

            def _find_placeholder(slide, types: set[int]) -> object | None:  # noqa: ANN001
                for shape in _iter_shapes(slide.shapes):
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    ph_type = _shape_ph_type(shape)
                    if ph_type is None:
                        continue
                    if ph_type in types:
                        return shape
                return None

            def _paragraph_font_score(paragraph) -> float:  # noqa: ANN001
                best = 0.0
                for run in getattr(paragraph, "runs", []) or []:
                    try:
                        if run.font.size is not None:
                            best = max(best, float(run.font.size.pt))
                    except Exception:
                        continue
                return best

            def _best_run_by_font(paragraph) -> object | None:  # noqa: ANN001
                runs = list(getattr(paragraph, "runs", []) or [])
                if not runs:
                    return None
                best_run = runs[0]
                best_score = -1.0
                for run in runs:
                    score = -1.0
                    try:
                        if run.font.size is not None:
                            score = float(run.font.size.pt)
                    except Exception:
                        score = -1.0
                    if score > best_score:
                        best_score = score
                        best_run = run
                return best_run

            def _copy_run_font(src_run, dest_run) -> None:  # noqa: ANN001
                try:
                    name = getattr(getattr(src_run, "font", None), "name", None)
                    if name:
                        dest_run.font.name = name
                except Exception:
                    pass
                try:
                    size = getattr(getattr(src_run, "font", None), "size", None)
                    if size is not None:
                        dest_run.font.size = size
                except Exception:
                    pass
                for attr in ("bold", "italic", "underline"):
                    try:
                        value = getattr(getattr(src_run, "font", None), attr, None)
                        if value is not None:
                            setattr(dest_run.font, attr, value)
                    except Exception:
                        continue
                try:
                    src_color = getattr(getattr(getattr(src_run, "font", None), "color", None), "rgb", None)
                    if src_color is not None:
                        dest_run.font.color.rgb = src_color
                except Exception:
                    pass

            def _set_paragraph_text_preserve(paragraph, value: str) -> None:  # noqa: ANN001
                runs = list(getattr(paragraph, "runs", []) or [])
                if not runs:
                    try:
                        paragraph.text = str(value or "")
                    except Exception:
                        return
                    return

                best_run = runs[0]
                best_score = -1.0
                for run in runs:
                    score = -1.0
                    try:
                        if run.font.size is not None:
                            score = float(run.font.size.pt)
                    except Exception:
                        score = -1.0
                    if score > best_score:
                        best_score = score
                        best_run = run

                try:
                    best_run.text = str(value or "")
                except Exception:
                    return
                for run in runs:
                    if run is best_run:
                        continue
                    try:
                        run.text = ""
                    except Exception:
                        continue

            def _set_shape_text_preserve(shape, value: str) -> None:  # noqa: ANN001
                if not getattr(shape, "has_text_frame", False):
                    return
                tf = shape.text_frame
                lines = str(value or "").splitlines()
                if not lines:
                    lines = [""]

                paragraphs = list(getattr(tf, "paragraphs", []) or [])
                if not paragraphs:
                    _set_shape_text(shape, "\n".join(lines))
                    return

                if len(lines) == 1 and len(paragraphs) > 1:
                    best_idx = 0
                    best_score = -1.0
                    for idx, p in enumerate(paragraphs):
                        score = _paragraph_font_score(p)
                        if score > best_score:
                            best_score = score
                            best_idx = idx
                    style_run = _best_run_by_font(paragraphs[best_idx])
                    dest = paragraphs[0]
                    _set_paragraph_text_preserve(dest, lines[0])
                    dest_run = _best_run_by_font(dest)
                    if style_run is not None and dest_run is not None:
                        _copy_run_font(style_run, dest_run)
                    for idx, p in enumerate(paragraphs):
                        if idx == 0:
                            continue
                        _set_paragraph_text_preserve(p, "")
                    return

                while len(paragraphs) < len(lines):
                    try:
                        paragraphs.append(tf.add_paragraph())
                    except Exception:
                        break
                for idx, line in enumerate(lines):
                    if idx >= len(paragraphs):
                        break
                    _set_paragraph_text_preserve(paragraphs[idx], line)
                for idx in range(len(lines), len(paragraphs)):
                    _set_paragraph_text_preserve(paragraphs[idx], "")

            def _clear_large_text(slide, keep_shape_ids: set[int]) -> None:  # noqa: ANN001
                for shape in _iter_shapes(slide.shapes):
                    try:
                        shape_id = int(getattr(shape, "shape_id", 0) or 0)
                    except Exception:
                        shape_id = 0
                    if shape_id and shape_id in keep_shape_ids:
                        continue
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    value = _shape_text(shape).strip()
                    if not value:
                        continue
                    ph_type = _shape_ph_type(shape)
                    if ph_type is not None:
                        if ph_type in {int(PP_PLACEHOLDER.FOOTER), int(PP_PLACEHOLDER.SLIDE_NUMBER)}:
                            continue
                        _set_shape_text(shape, "")
                        continue

                    # Keep common footer elements and small copyright lines.
                    if "©" in value or "copyright" in value.lower():
                        continue
                    try:
                        if int(shape.top) >= int(slide_height * 0.92):
                            continue
                    except Exception:
                        pass

                    _set_shape_text(shape, "")

            def _strip_non_background_pictures(slide) -> None:  # noqa: ANN001
                def _group_contains_picture(group) -> bool:  # noqa: ANN001
                    try:
                        for child in group.shapes:
                            st = getattr(child, "shape_type", None)
                            if st == MSO_SHAPE_TYPE.PICTURE:
                                return True
                            if st == MSO_SHAPE_TYPE.GROUP and _group_contains_picture(child):
                                return True
                    except Exception:
                        return False
                    return False

                total_area = max(1, int(slide_width) * int(slide_height))
                for shape in list(slide.shapes):
                    shape_type = getattr(shape, "shape_type", None)
                    is_picture = shape_type == MSO_SHAPE_TYPE.PICTURE
                    if shape_type == MSO_SHAPE_TYPE.GROUP and _group_contains_picture(shape):
                        # Some templates wrap pictures inside groups; remove the whole group to avoid leftovers.
                        is_picture = True
                    if not is_picture:
                        continue
                    if _is_background_shape(shape, slide_width, slide_height):
                        continue
                    try:
                        area_ratio = float(_shape_area(shape)) / float(total_area)
                        top = int(shape.top)
                    except Exception:
                        area_ratio = 1.0
                        top = 0
                    # Keep tiny header logos/icons; drop the rest.
                    if area_ratio <= 0.02 and top <= int(slide_height * 0.22):
                        continue
                    _remove_shape(shape)

            def _fill_cover_slide_inplace(slide) -> None:  # noqa: ANN001
                text_shapes = [shape for shape in _iter_shapes(slide.shapes) if getattr(shape, "has_text_frame", False)]
                title_shape = _find_placeholder(slide, {int(PP_PLACEHOLDER.TITLE), int(PP_PLACEHOLDER.CENTER_TITLE)})
                if title_shape is None:
                    if not text_shapes:
                        return
                    title_shape = max(text_shapes, key=_shape_font_score)
                try:
                    title_shape_id = int(getattr(title_shape, "shape_id", 0) or 0)
                except Exception:
                    title_shape_id = 0

                subtitle_shape = _find_placeholder(slide, {int(PP_PLACEHOLDER.SUBTITLE)})
                date_shape = _find_placeholder(slide, {int(PP_PLACEHOLDER.DATE)})
                subtitle_shape_id = 0
                if subtitle_shape is not None:
                    try:
                        subtitle_shape_id = int(getattr(subtitle_shape, "shape_id", 0) or 0)
                    except Exception:
                        subtitle_shape_id = 0

                title_lines = _cover_title_lines(normalized_title)
                if subtitle_shape is None:
                    for shape in text_shapes:
                        try:
                            sid = int(getattr(shape, "shape_id", 0) or 0)
                        except Exception:
                            sid = 0
                        if sid and title_shape_id and sid == title_shape_id:
                            continue
                        value = _shape_text(shape).strip()
                        if not value:
                            continue
                        if "©" in value or "copyright" in value.lower():
                            continue
                        try:
                            if int(shape.top) <= int(title_shape.top):
                                continue
                        except Exception:
                            continue
                        if subtitle_shape is None or int(getattr(shape, "width", 0) or 0) > int(
                            getattr(subtitle_shape, "width", 0) or 0
                        ):
                            subtitle_shape = shape
                            subtitle_shape_id = sid

                if date_shape is None:
                    for shape in text_shapes:
                        try:
                            sid = int(getattr(shape, "shape_id", 0) or 0)
                        except Exception:
                            sid = 0
                        if sid and title_shape_id and sid == title_shape_id:
                            continue
                        if sid and subtitle_shape_id and sid == subtitle_shape_id:
                            continue
                        value = _shape_text(shape)
                        if not value:
                            continue
                        if re.search(r"\d{4}[-/.]\d{1,2}[-/.]\d{1,2}", value):
                            date_shape = shape
                            break

                _set_shape_text_preserve(title_shape, "\n".join(title_lines))
                if subtitle_shape is not None:
                    _set_shape_text_preserve(subtitle_shape, cover_subtitle)
                if date_shape is not None:
                    _set_shape_text_preserve(date_shape, datetime.now().strftime("%Y-%m-%d"))

                keep_shape_ids: set[int] = set()
                for shape in (title_shape, subtitle_shape, date_shape):
                    if shape is None:
                        continue
                    try:
                        shape_id = int(getattr(shape, "shape_id", 0) or 0)
                    except Exception:
                        shape_id = 0
                    if shape_id:
                        keep_shape_ids.add(shape_id)
                if keep_shape_ids:
                    _clear_large_text(slide, keep_shape_ids)

            def _fill_body_bullets_inplace(shape, bullets: list[str]) -> None:  # noqa: ANN001
                if not getattr(shape, "has_text_frame", False):
                    return
                tf = shape.text_frame
                tf.word_wrap = True
                items = bullets or ["结合业务场景补充关键结论与下一步行动。"]
                items = [str(x).strip() for x in items if str(x).strip()]
                if not items:
                    items = ["结合业务场景补充关键结论与下一步行动。"]

                ph_type = _shape_ph_type(shape)
                assume_bullets = ph_type in {int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)}
                prefix = "" if assume_bullets else "• "

                paragraphs = list(tf.paragraphs or [])
                if not paragraphs:
                    try:
                        tf.text = ""
                    except Exception:
                        return
                    paragraphs = list(tf.paragraphs or [])
                for idx, item in enumerate(items[:8]):
                    text_value = f"{prefix}{item}" if prefix else item
                    if idx < len(paragraphs):
                        _set_paragraph_text_preserve(paragraphs[idx], text_value)
                        try:
                            paragraphs[idx].level = 0
                        except Exception:
                            pass
                    else:
                        p = tf.add_paragraph()
                        _set_paragraph_text_preserve(p, text_value)
                        try:
                            p.level = 0
                        except Exception:
                            pass
                for idx in range(len(items[:8]), len(paragraphs)):
                    _set_paragraph_text_preserve(paragraphs[idx], "")

            def _fill_content_slide_inplace(slide, slide_title: str, bullets: list[str]) -> None:  # noqa: ANN001
                title_shape = _find_placeholder(slide, {int(PP_PLACEHOLDER.TITLE), int(PP_PLACEHOLDER.CENTER_TITLE)})
                body_shape = _find_placeholder(slide, {int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)})

                def _is_title_like_body_placeholder(shape) -> bool:  # noqa: ANN001
                    if shape is None:
                        return False
                    ph_type = _shape_ph_type(shape)
                    if ph_type not in {int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)}:
                        return False
                    try:
                        top = int(shape.top)
                        height = int(shape.height)
                        width = int(shape.width)
                    except Exception:
                        return False
                    # Some templates (including JetLinks) use a BODY placeholder as a top title band.
                    if top <= int(Inches(1.15)) and height <= int(Inches(1.05)) and width >= int(slide_width * 0.6):
                        return True
                    return False

                if title_shape is None and body_shape is not None and _is_title_like_body_placeholder(body_shape):
                    title_shape = body_shape
                    body_shape = None

                if title_shape is None:
                    title_shape = _pick_title_shape(slide)
                if body_shape is None:
                    body_shape = _pick_body_shape(slide, title_shape)

                if title_shape is not None and body_shape is not None and getattr(title_shape, "shape_id", None) == getattr(body_shape, "shape_id", None):
                    title_shape = None

                if title_shape is None:
                    title_shape = slide.shapes.add_textbox(Inches(0.9), Inches(0.50), Inches(11.2), Inches(0.62))

                if body_shape is None:
                    body_shape = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.2), Inches(4.9))

                keep_shape_ids: set[int] = set()
                if title_shape is not None:
                    _set_shape_text_preserve(title_shape, _normalize_text(slide_title, max_chars=34, fallback="—"))
                    try:
                        shape_id = int(getattr(title_shape, "shape_id", 0) or 0)
                    except Exception:
                        shape_id = 0
                    if shape_id:
                        keep_shape_ids.add(shape_id)
                if body_shape is not None:
                    _fill_body_bullets_inplace(body_shape, bullets)
                    try:
                        shape_id = int(getattr(body_shape, "shape_id", 0) or 0)
                    except Exception:
                        shape_id = 0
                    if shape_id:
                        keep_shape_ids.add(shape_id)
                if keep_shape_ids:
                    _clear_large_text(slide, keep_shape_ids)

            keep_images = _effective_keep_images()
            mode = _effective_template_mode()
            template_indices_used: list[int] | None = None

            if mode in {"reuse", "inplace", "preserve"}:
                # Reuse template slides in-place (no cloning) to preserve design fidelity and avoid broken relations.
                if len(prs_t.slides) < needed + 1:
                    raise ValueError("template has insufficient slides for reuse mode")

                # Optional: user-specified content slide indices (1-based). If omitted, fall back to env config,
                # and finally auto-pick the cleanest content slides.
                configured_indices = _normalize_template_indices(template_content_indices, len(prs_t.slides))
                raw_indices = os.getenv("AISTAFF_PPT_TEMPLATE_CONTENT_INDICES") or ""
                parsed_indices = _parse_template_indices(raw_indices, len(prs_t.slides))
                parsed_indices = [idx for idx in parsed_indices if idx != 1]

                cover_slide = prs_t.slides[0]
                content_slides = list(prs_t.slides)[1 : 1 + needed]

                selected: list[int] = []
                if configured_indices:
                    selected = configured_indices[:needed]
                elif parsed_indices:
                    selected = parsed_indices[:needed]

                # If indices are missing/insufficient, auto-pick the "cleanest" slides from the template
                # to avoid reusing complex business pages (lots of pictures/groups/textboxes).
                if len(selected) < needed:
                    slide_area = max(1, int(slide_width) * int(slide_height))
                    candidates: list[tuple[tuple[int, ...], int]] = []
                    for idx, slide in enumerate(list(prs_t.slides), start=1):
                        if idx == 1:
                            continue
                        shapes = list(slide.shapes)
                        shape_count = len(shapes)
                        text_count = sum(1 for s in shapes if getattr(s, "has_text_frame", False))
                        nonempty_text = sum(
                            1
                            for s in shapes
                            if getattr(s, "has_text_frame", False) and str(getattr(s, "text", "") or "").strip()
                        )
                        group_count = sum(1 for s in shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.GROUP)
                        pic_count = sum(1 for s in shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE)
                        if text_count == 0:
                            continue
                        pic_area = 0
                        max_pic_ratio = 0.0
                        for s in shapes:
                            if getattr(s, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                                continue
                            if _is_background_shape(s, slide_width, slide_height):
                                continue
                            area = _shape_area(s)
                            if area <= 0:
                                continue
                            pic_area += area
                            ratio = float(area) / float(slide_area)
                            if ratio > max_pic_ratio:
                                max_pic_ratio = ratio
                        pic_area_scaled = int(round((float(pic_area) / float(slide_area)) * 1000.0)) if slide_area else 0
                        max_pic_scaled = int(round(max_pic_ratio * 1000.0))
                        body_ph = sum(
                            1
                            for s in shapes
                            if _shape_ph_type(s) in {int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)}
                        )

                        # Estimate available content area below the header placeholder. We prefer templates that provide
                        # a large text box region for bullets, avoiding "title-only" or overly specialized pages.
                        header_shape = None
                        for s in shapes:
                            ph = _shape_ph_type(s)
                            if ph not in {
                                int(PP_PLACEHOLDER.TITLE),
                                int(PP_PLACEHOLDER.CENTER_TITLE),
                                int(PP_PLACEHOLDER.BODY),
                                int(PP_PLACEHOLDER.OBJECT),
                            }:
                                continue
                            if header_shape is None or int(getattr(s, "top", 0) or 0) < int(getattr(header_shape, "top", 0) or 0):
                                header_shape = s

                        header_top = int(getattr(header_shape, "top", 0) or 0) if header_shape is not None else 0
                        header_id = int(getattr(header_shape, "shape_id", 0) or 0) if header_shape is not None else 0
                        min_top = header_top + int(Inches(0.80))

                        max_content_area = 0
                        for s in shapes:
                            if not getattr(s, "has_text_frame", False):
                                continue
                            try:
                                sid = int(getattr(s, "shape_id", 0) or 0)
                            except Exception:
                                sid = 0
                            if header_id and sid == header_id:
                                continue
                            ph = _shape_ph_type(s)
                            if ph in {int(PP_PLACEHOLDER.FOOTER), int(PP_PLACEHOLDER.SLIDE_NUMBER), int(PP_PLACEHOLDER.DATE)}:
                                continue
                            try:
                                top = int(s.top)
                            except Exception:
                                top = 0
                            if top < min_top:
                                continue
                            if top >= int(slide_height * 0.92):
                                continue
                            area = _shape_area(s)
                            if area > max_content_area:
                                max_content_area = area

                        content_scaled = int(round((float(max_content_area) / float(slide_area)) * 1000.0))

                        # Buckets (smaller is better): prefer slides that look like reusable "content layouts".
                        # Note: some templates have no placeholders; in that case we still pick the least complex slides.
                        if body_ph > 0 and content_scaled >= 150 and group_count == 0 and pic_count <= 3 and shape_count <= 40:
                            bucket = 0
                        elif body_ph > 0 and content_scaled >= 100 and group_count <= 1 and pic_count <= 4 and shape_count <= 60:
                            bucket = 1
                        elif body_ph > 0 and content_scaled >= 70:
                            bucket = 2
                        elif body_ph > 0:
                            bucket = 3
                        elif content_scaled >= 120 and group_count == 0 and pic_count <= 4 and shape_count <= 60:
                            bucket = 4
                        else:
                            bucket = 5
                        # Extra penalty: slides with a large non-background picture are often image-centric
                        # layouts (little space for bullets), even if they look "simple" by shape count.
                        if max_pic_scaled >= 260:
                            bucket += 2
                        elif max_pic_scaled >= 180:
                            bucket += 1

                        # Penalize slides that contain lots of pre-filled text (likely business pages).
                        # Also avoid layouts that reserve big areas for non-background pictures (image-centric pages).
                        score = (
                            bucket,
                            max_pic_scaled,
                            pic_area_scaled,
                            -content_scaled,
                            nonempty_text,
                            shape_count,
                            pic_count,
                            group_count,
                            text_count,
                            idx,
                        )
                        candidates.append((score, idx))
                    candidates.sort(key=lambda item: item[0])
                    for _, idx in candidates:
                        if idx not in selected:
                            selected.append(idx)
                        if len(selected) >= needed:
                            break
                template_indices_used = selected[:needed] if selected else None

                # If we have enough configured/auto-picked indices, keep exactly those slides (1-based indices),
                # and drop everything else. This avoids reusing "wrong" layouts that contain unwanted images/text.
                if len(selected) >= needed:
                    selected = selected[:needed]
                    keep_set = {0} | {max(0, int(i) - 1) for i in selected}
                    for idx in range(len(prs_t.slides) - 1, -1, -1):
                        if idx not in keep_set:
                            _remove_slide(prs_t, idx)
                    cover_slide = prs_t.slides[0]
                    content_slides = list(prs_t.slides)[1:]

                # Drop extra template slides so output is deterministic.
                for idx in range(len(prs_t.slides) - 1, 0 + needed, -1):
                    _remove_slide(prs_t, idx)

                if not keep_images:
                    _strip_non_background_pictures(cover_slide)
                    for slide in content_slides:
                        _strip_non_background_pictures(slide)

                _fill_cover_slide_inplace(cover_slide)
                for slide, item in zip(content_slides, normalized_slides):
                    _fill_content_slide_inplace(
                        slide,
                        str(item.get("title") or ""),
                        [str(x) for x in (item.get("bullets") or [])],
                    )
            else:
                raw_indices = os.getenv("AISTAFF_PPT_TEMPLATE_CONTENT_INDICES") or ""
                parsed_indices = _parse_template_indices(raw_indices, len(prs_t.slides))
                parsed_indices = [idx for idx in parsed_indices if idx != 1]
                if not parsed_indices:
                    base_idx = _pick_base_content_slide_index(prs_t)
                    parsed_indices = [base_idx + 1]
                template_indices_used = parsed_indices[:]

                base_indices = [idx - 1 for idx in parsed_indices if 0 <= idx - 1 < len(prs_t.slides)]
                if not base_indices:
                    base_indices = [max(0, _pick_base_content_slide_index(prs_t))]

                orig_slides = list(prs_t.slides)
                content_slides = []
                for i in range(needed):
                    source = orig_slides[base_indices[i % len(base_indices)]]
                    new_slide = _duplicate_slide(prs_t, source, blank_layout)
                    _cleanup_slide(new_slide, slide_width, slide_height, remove_pictures=not keep_images)
                    content_slides.append(new_slide)

                for idx in range(len(orig_slides) - 1, 0, -1):
                    _remove_slide(prs_t, idx)

                cover_slide = prs_t.slides[0]
                _cleanup_slide(cover_slide, slide_width, slide_height, remove_pictures=False)
                _fill_cover_slide(cover_slide)

                for slide, item in zip(content_slides, normalized_slides):
                    _fill_clean_content_slide(
                        slide,
                        str(item.get("title") or ""),
                        [str(x) for x in (item.get("bullets") or [])],
                    )

            filename = f"{uuid4().hex}.pptx"
            out_path = (self._settings.outputs_dir / filename).resolve()
            prs_t.save(str(out_path))

            result = _build_result(out_path)
            result.update(
                {
                    "render_mode": "template",
                    "template_used": candidate.name,
                    "template_mode": mode,
                    "template_keep_images": keep_images,
                    "template_content_indices": template_indices_used or [],
                }
            )
            return result

        candidates: list[Path] = []
        if template_override is not None:
            candidates.append(template_override)
        candidates.extend(self._ppt_template_candidates(normalized_style))
        seen_candidates: set[str] = set()
        filtered: list[Path] = []
        for c in candidates:
            key = str(c)
            if key in seen_candidates:
                continue
            seen_candidates.add(key)
            filtered.append(c)

        for candidate in filtered:
            if not candidate.exists() or not candidate.is_file():
                continue
            try:
                return _render_with_template(candidate)
            except Exception as e:
                template_failures.append(
                    {
                        "template": candidate.name,
                        "error": _normalize_text(str(e), max_chars=260, fallback="template render failed"),
                    }
                )
                if strict_template:
                    raise
                continue

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        font_name = str(os.getenv("AISTAFF_PPT_FONT") or "").strip() or "微软雅黑"
        palette = style_profiles.get(effective_style, style_profiles["modern_blue"])
        navy = _rgb(palette["cover_from"])
        deep_blue = _rgb(palette["cover_to"])
        primary = _rgb(palette["primary"])
        accent = _rgb(palette["accent"])
        accent2 = _rgb(palette["accent2"])
        text_color = _rgb(palette["text"])
        muted = _rgb(palette["muted"])
        border = _rgb(palette["border"])
        bg = _rgb(palette["bg"])
        cover_text = _rgb(palette["cover_text"])
        surface = _rgb(palette["surface"])
        white = _rgb("FFFFFF")
        chip_bg = _rgb(palette["chip_bg"])
        success = _rgb(palette["success"])

        def _blend(a: RGBColor, b: RGBColor, alpha: float) -> RGBColor:
            ratio = max(0.0, min(1.0, float(alpha)))
            return RGBColor(
                int(round((a[0] * ratio) + (b[0] * (1 - ratio)))),
                int(round((a[1] * ratio) + (b[1] * (1 - ratio)))),
                int(round((a[2] * ratio) + (b[2] * (1 - ratio)))),
            )

        blank = prs.slide_layouts[6]
        total_pages = len(normalized_slides) + 1
        footer_title = _normalize_text(normalized_title, max_chars=34, fallback="演示文稿")

        def add_shape_fill(slide, shape_type: MSO_SHAPE, left: float, top: float, width: float, height: float, color: RGBColor, transparency: float = 0.0):  # noqa: ANN001
            shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            if 0 < transparency < 1:
                shape.fill.transparency = transparency
            shape.line.fill.background()
            return shape

        def add_cover_background(slide) -> None:  # noqa: ANN001
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg_shape.fill.gradient()
            bg_shape.fill.gradient_angle = 38
            stops = bg_shape.fill.gradient_stops
            stops[0].position = 0.0
            stops[0].color.rgb = navy
            stops[1].position = 1.0
            stops[1].color.rgb = deep_blue
            bg_shape.line.fill.background()

            orb1 = add_shape_fill(slide, MSO_SHAPE.OVAL, 9.6, -0.9, 4.7, 4.7, accent2, transparency=0.18)
            orb1.fill.gradient()
            orb1.fill.gradient_angle = 45
            orb1_stops = orb1.fill.gradient_stops
            orb1_stops[0].position = 0.0
            orb1_stops[0].color.rgb = accent2
            orb1_stops[1].position = 1.0
            orb1_stops[1].color.rgb = accent
            orb1.fill.transparency = 0.2

            add_shape_fill(slide, MSO_SHAPE.OVAL, 10.8, 5.2, 2.7, 2.7, accent, transparency=0.55)
            add_shape_fill(slide, MSO_SHAPE.OVAL, -0.9, 5.7, 2.6, 2.6, accent2, transparency=0.5)

        def add_cover_meta(slide) -> None:  # noqa: ANN001
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.72), Inches(3.15), Inches(0.54))
            badge.fill.solid()
            badge.fill.fore_color.rgb = white
            badge.fill.transparency = 0.05
            badge.line.fill.background()

            tx = slide.shapes.add_textbox(Inches(1.08), Inches(0.83), Inches(2.8), Inches(0.34))
            tf = tx.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "AISTAFF · 演示模板"
            p.alignment = PP_ALIGN.LEFT
            p.font.name = font_name
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = accent

        def add_cover_title(slide, title_lines: list[str]) -> None:  # noqa: ANN001
            title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.85), prs.slide_width - Inches(1.9), Inches(2.65))
            tf = title_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP

            full_text_len = sum(len(x) for x in title_lines)
            title_size = _cover_title_size(full_text_len)

            for idx, line in enumerate(title_lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)
                p.line_spacing = 1.12
                p.font.name = font_name
                p.font.bold = True
                p.font.size = title_size
                p.font.color.rgb = cover_text

        def add_cover_subtitle(slide, subtitle: str) -> None:  # noqa: ANN001
            subtitle_box = slide.shapes.add_textbox(Inches(0.92), Inches(4.65), prs.slide_width - Inches(2.1), Inches(0.92))
            subtitle_tf = subtitle_box.text_frame
            subtitle_tf.clear()
            subtitle_tf.word_wrap = True
            p = subtitle_tf.paragraphs[0]
            p.text = subtitle
            p.font.name = font_name
            p.font.size = Pt(18)
            p.font.color.rgb = cover_text
            p.line_spacing = 1.18

            meta = subtitle_tf.add_paragraph()
            meta.text = f"生成日期：{datetime.now().strftime('%Y-%m-%d')}  ·  共 {total_pages} 页"
            meta.font.name = font_name
            meta.font.size = Pt(12)
            meta.font.color.rgb = cover_text

        def add_content_background(slide, page_no: int) -> None:  # noqa: ANN001
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = bg
            bg_shape.line.fill.background()

            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.13))
            top_bar.fill.gradient()
            top_bar.fill.gradient_angle = 0
            stops = top_bar.fill.gradient_stops
            stops[0].position = 0.0
            stops[0].color.rgb = accent
            stops[1].position = 1.0
            stops[1].color.rgb = accent2
            top_bar.line.fill.background()

            add_shape_fill(slide, MSO_SHAPE.OVAL, 11.2, -0.55, 2.4, 2.4, accent2, transparency=0.82)
            add_shape_fill(slide, MSO_SHAPE.OVAL, -0.7, 6.2, 2.0, 2.0, accent, transparency=0.88)

            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                prs.slide_width - Inches(2.35),
                Inches(0.44),
                Inches(1.55),
                Inches(0.42),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = chip_bg
            badge.line.fill.background()

            badge_text = slide.shapes.add_textbox(prs.slide_width - Inches(2.25), Inches(0.50), Inches(1.35), Inches(0.30))
            badge_tf = badge_text.text_frame
            badge_tf.clear()
            badge_p = badge_tf.paragraphs[0]
            badge_p.text = f"P{page_no:02d}"
            badge_p.alignment = PP_ALIGN.CENTER
            badge_p.font.name = font_name
            badge_p.font.size = Pt(11)
            badge_p.font.bold = True
            badge_p.font.color.rgb = accent

        def add_content_title(slide, value: str) -> None:  # noqa: ANN001
            title_text = _normalize_text(value, max_chars=34, fallback="—")
            title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.52), prs.slide_width - Inches(3.4), Inches(0.82))
            tf = title_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = PP_ALIGN.LEFT
            p.font.name = font_name
            p.font.bold = True
            p.font.size = _title_size(len(title_text))
            p.font.color.rgb = primary

            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(1.24), Inches(1.5), Inches(0.06))
            line.fill.solid()
            line.fill.fore_color.rgb = accent
            line.line.fill.background()

        def add_footer(slide, page_no: int) -> None:  # noqa: ANN001
            left = slide.shapes.add_textbox(Inches(0.92), prs.slide_height - Inches(0.42), Inches(9.1), Inches(0.24))
            left_tf = left.text_frame
            left_tf.clear()
            left_p = left_tf.paragraphs[0]
            left_p.text = footer_title
            left_p.alignment = PP_ALIGN.LEFT
            left_p.font.name = font_name
            left_p.font.size = Pt(10.5)
            left_p.font.color.rgb = muted

            right = slide.shapes.add_textbox(prs.slide_width - Inches(1.65), prs.slide_height - Inches(0.42), Inches(1.45), Inches(0.24))
            right_tf = right.text_frame
            right_tf.clear()
            right_p = right_tf.paragraphs[0]
            right_p.text = f"{page_no}/{total_pages}"
            right_p.alignment = PP_ALIGN.RIGHT
            right_p.font.name = font_name
            right_p.font.size = Pt(10.5)
            right_p.font.color.rgb = muted

        def add_bullet_block(textbox, items: list[str], *, font_size: Pt, bullet_color: RGBColor = accent):  # noqa: ANN001
            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP

            for idx, item in enumerate(items):
                paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.line_spacing = 1.18
                paragraph.space_after = Pt(6)

                mark = paragraph.add_run()
                mark.text = "• "
                mark.font.name = font_name
                mark.font.size = Pt(max(10, int(font_size.pt) - 1))
                mark.font.bold = False
                mark.font.color.rgb = bullet_color

                run = paragraph.add_run()
                run.text = item
                run.font.name = font_name
                run.font.size = font_size
                run.font.color.rgb = text_color

        cover = prs.slides.add_slide(blank)
        add_cover_background(cover)
        add_cover_meta(cover)
        add_cover_title(cover, _cover_title_lines(normalized_title))
        add_cover_subtitle(cover, cover_subtitle)

        for page_index, slide_data in enumerate(normalized_slides, 1):
            content_slide = prs.slides.add_slide(blank)
            slide_title = str(slide_data.get("title") or "").strip() or "—"
            bullets = [str(item) for item in (slide_data.get("bullets") or []) if str(item).strip()]

            page_no = page_index + 1
            add_content_background(content_slide, page_no=page_no)
            add_content_title(content_slide, slide_title)

            max_len = max((len(item) for item in bullets), default=0)
            auto_layout = "single_column"
            if len(bullets) <= 4 and max_len <= 36:
                auto_layout = "focus"
            elif 3 <= len(bullets) <= 6 and max_len <= 30:
                auto_layout = "cards"
            elif len(bullets) >= 7 or (len(bullets) >= 5 and max_len >= 34):
                auto_layout = "two_column"

            chosen_layout = normalized_layout_mode if normalized_layout_mode != "auto" else auto_layout

            card_left = 0.9
            card_top = 1.46
            card_width = 11.52
            card_height = 5.43

            card = content_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(card_left),
                Inches(card_top),
                Inches(card_width),
                Inches(card_height),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = surface
            card.line.color.rgb = border

            if chosen_layout == "focus":
                key_fill = _blend(chip_bg, surface, 0.75) if effective_style == "dark_tech" else _blend(chip_bg, surface, 0.35)
                key_box = content_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(card_left + 0.34),
                    Inches(card_top + 0.32),
                    Inches(6.5),
                    Inches(4.72),
                )
                key_box.fill.solid()
                key_box.fill.fore_color.rgb = key_fill
                key_box.line.fill.background()

                key_title = content_slide.shapes.add_textbox(Inches(card_left + 0.58), Inches(card_top + 0.56), Inches(6.0), Inches(0.4))
                key_title_tf = key_title.text_frame
                key_title_tf.clear()
                key_title_p = key_title_tf.paragraphs[0]
                key_title_p.text = "关键结论"
                key_title_p.font.name = font_name
                key_title_p.font.bold = True
                key_title_p.font.size = Pt(14)
                key_title_p.font.color.rgb = accent

                key_body = content_slide.shapes.add_textbox(Inches(card_left + 0.58), Inches(card_top + 1.02), Inches(6.02), Inches(3.86))
                key_tf = key_body.text_frame
                key_tf.clear()
                key_tf.word_wrap = True
                key_tf.vertical_anchor = MSO_ANCHOR.TOP
                key_p = key_tf.paragraphs[0]
                key_p.text = bullets[0]
                key_p.alignment = PP_ALIGN.LEFT
                key_p.line_spacing = 1.22
                key_p.font.name = font_name
                key_p.font.bold = True
                key_p.font.size = Pt(25 if len(bullets[0]) <= 20 else 21)
                key_p.font.color.rgb = primary

                rest = bullets[1:] if len(bullets) > 1 else ["结合现场情况补充案例、数据与行动计划。"]
                side_title = content_slide.shapes.add_textbox(Inches(card_left + 7.02), Inches(card_top + 0.56), Inches(3.95), Inches(0.34))
                side_title_tf = side_title.text_frame
                side_title_tf.clear()
                side_title_p = side_title_tf.paragraphs[0]
                side_title_p.text = "行动要点"
                side_title_p.font.name = font_name
                side_title_p.font.bold = True
                side_title_p.font.size = Pt(13)
                side_title_p.font.color.rgb = accent

                side_body = content_slide.shapes.add_textbox(Inches(card_left + 7.00), Inches(card_top + 0.94), Inches(4.0), Inches(4.6))
                add_bullet_block(side_body, rest[:5], font_size=Pt(17), bullet_color=success)
            elif chosen_layout == "cards":
                card_items = bullets[:6] if bullets else ["结合业务场景补充要点。"]
                columns = 3 if len(card_items) >= 5 else 2
                rows = (len(card_items) + columns - 1) // columns
                inner_left = card_left + 0.35
                inner_top = card_top + 0.35
                inner_width = card_width - 0.70
                inner_height = card_height - 0.70
                gutter_x = 0.22
                gutter_y = 0.22
                box_width = (inner_width - (columns - 1) * gutter_x) / columns
                box_height = (inner_height - (rows - 1) * gutter_y) / max(1, rows)

                for idx, item in enumerate(card_items):
                    r = idx // columns
                    c = idx % columns
                    x = inner_left + c * (box_width + gutter_x)
                    y = inner_top + r * (box_height + gutter_y)
                    tip_fill = _blend(chip_bg, surface, 0.82) if effective_style == "dark_tech" else _blend(chip_bg, surface, 0.22)

                    tip = content_slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x),
                        Inches(y),
                        Inches(box_width),
                        Inches(box_height),
                    )
                    tip.fill.solid()
                    tip.fill.fore_color.rgb = tip_fill
                    tip.line.color.rgb = border

                    icon = content_slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(x + 0.16),
                        Inches(y + 0.16),
                        Inches(0.18),
                        Inches(0.18),
                    )
                    icon.fill.solid()
                    icon.fill.fore_color.rgb = accent
                    icon.line.fill.background()

                    tip_box = content_slide.shapes.add_textbox(
                        Inches(x + 0.40),
                        Inches(y + 0.12),
                        Inches(max(0.8, box_width - 0.55)),
                        Inches(max(0.5, box_height - 0.24)),
                    )
                    tip_tf = tip_box.text_frame
                    tip_tf.clear()
                    tip_tf.word_wrap = True
                    tip_p = tip_tf.paragraphs[0]
                    tip_p.text = item
                    tip_p.alignment = PP_ALIGN.LEFT
                    tip_p.font.name = font_name
                    tip_p.font.size = Pt(14 if len(item) <= 20 else 12.5)
                    tip_p.font.bold = len(item) <= 18
                    tip_p.font.color.rgb = text_color
            else:
                body_font = _font_size_for_bullets(bullets)
                if chosen_layout == "two_column":
                    inner_left = card_left + 0.34
                    inner_top = card_top + 0.30
                    inner_width = card_width - 0.68
                    inner_height = card_height - 0.60
                    gutter = 0.30
                    column_width = (inner_width - gutter) / 2
                    left_items, right_items = _split_bullets(bullets)

                    divider = content_slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(inner_left + column_width + (gutter / 2) - 0.01),
                        Inches(inner_top + 0.10),
                        Inches(0.02),
                        Inches(inner_height - 0.20),
                    )
                    divider.fill.solid()
                    divider.fill.fore_color.rgb = border
                    divider.line.fill.background()

                    left_box = content_slide.shapes.add_textbox(Inches(inner_left), Inches(inner_top), Inches(column_width), Inches(inner_height))
                    right_box = content_slide.shapes.add_textbox(
                        Inches(inner_left + column_width + gutter),
                        Inches(inner_top),
                        Inches(column_width),
                        Inches(inner_height),
                    )

                    smaller = Pt(max(15, int(body_font.pt) - 1))
                    add_bullet_block(left_box, left_items, font_size=smaller)
                    add_bullet_block(right_box, right_items, font_size=smaller)
                else:
                    body = content_slide.shapes.add_textbox(
                        Inches(card_left + 0.42),
                        Inches(card_top + 0.34),
                        Inches(card_width - 0.84),
                        Inches(card_height - 0.68),
                    )
                    add_bullet_block(body, bullets, font_size=body_font)

            add_footer(content_slide, page_no=page_no)

        filename = f"{uuid4().hex}.pptx"
        path = (self._settings.outputs_dir / filename).resolve()
        prs.save(str(path))
        result = _build_result(path)
        if template_failures:
            result["template_failures"] = template_failures
        result["render_mode"] = "programmatic"
        return result
